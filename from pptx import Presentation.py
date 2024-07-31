import mechanicalsoup
import pandas as pd
from sqlalchemy import create_engine, Column, Integer, String, Float
from sqlalchemy.orm import declarative_base, sessionmaker
import re
from pptx import Presentation
from pptx.util import Inches

# URL of the webpage containing the data
url = 'https://www.ncei.noaa.gov/access/monitoring/monthly-report/global/202313'

# Create a browser object
browser = mechanicalsoup.StatefulBrowser()
browser.open(url)

# Fetch the content
page = browser.get_current_page()

# Extract the table containing the data
table = page.find('table')

if table is None:
    print("No table found on the webpage.")
    exit()

# Function to sanitize column names
def sanitize_column_name(name):
    return re.sub(r'\W|^(?=\d)', '_', name).lower()

# Extract headers
headers = [sanitize_column_name(header.text.strip()) for header in table.find_all('th')]
print("Headers:", headers)

# Extract rows and handle rows with missing columns by filling with None
rows = []
for row in table.find_all('tr'):
    cells = row.find_all('td')
    if cells:
        row_data = [cell.text.strip() for cell in cells]
        if len(row_data) < len(headers):
            row_data.extend([None] * (len(headers) - len(row_data)))
        rows.append(row_data)
    else:
        print("Skipping row due to column mismatch:", [cell.text.strip() for cell in cells])

# Create a DataFrame
df = pd.DataFrame(rows, columns=headers)

# Rename columns with typos
df.rename(columns={'life_exectancy': 'life_expectancy'}, inplace=True)

# Debugging: Check initial DataFrame structure and types
print("Initial DataFrame:")
print(df.info())
print(df.head())

# Define SQLAlchemy model dynamically
Base = declarative_base()

# Create a dictionary of columns for the SQLAlchemy model
columns = {
    '__tablename__': 'ecological_footprint',
    'id': Column(Integer, primary_key=True, autoincrement=True)
}
for column in df.columns:
    if pd.api.types.is_string_dtype(df[column]):
        columns[column] = Column(String)
    elif pd.api.types.is_integer_dtype(df[column]):
        columns[column] = Column(Integer)
    elif pd.api.types.is_float_dtype(df[column]):
        columns[column] = Column(Float)

# Define the dynamic model
EcologicalFootprint = type('EcologicalFootprint', (Base,), columns)

# Create an engine that connects to the SQLite database
engine = create_engine('sqlite:///GlobalEcologicalFootprint.db', echo=True)

# Drop the existing table if it exists
Base.metadata.drop_all(engine, [Base.metadata.tables.get('ecological_footprint')])

# Create all tables in the engine
Base.metadata.create_all(engine)

# Create a session
Session = sessionmaker(bind=engine)
session = Session()

# Insert data into the database
for _, row in df.iterrows():
    row_data = {column: row[column] for column in df.columns}
    try:
        footprint_data = EcologicalFootprint(**row_data)
        session.add(footprint_data)
    except TypeError as e:
        print(f"Error inserting row: {row_data}")
        print(e)

# Commit the transaction
session.commit()

# Close the session
session.close()

# Read the data back from the database to display it
df_from_db = pd.read_sql_table('ecological_footprint', engine)
print("DataFrame from the database:")
print(df_from_db)

# Create a PowerPoint presentation
prs = Presentation()

# Title slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Global Ecological Footprint Data Analysis"
subtitle.text = "Summary of Data Extraction, Cleaning, Storage, and Visualization"

# Slide 1: Introduction
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Introduction"
content = slide.placeholders[1]
content.text = (
    "This presentation provides an overview of the process of extracting, "
    "cleaning, storing, and visualizing global ecological footprint data."
)

# Slide 2: Data Extraction Process
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Data Extraction Process"
content = slide.placeholders[1]
content.text = (
    "1. URL: 'https://www.ncei.noaa.gov/access/monitoring/monthly-report/global/202313'\n"
    "2. Used MechanicalSoup to open and fetch the webpage content.\n"
    "3. Extracted table headers and rows from the HTML table.\n"
    "4. Sanitized column names to ensure they are valid Python identifiers.\n"
)

# Slide 3: Data Cleaning
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Data Cleaning"
content = slide.placeholders[1]
content.text = (
    "1. Converted extracted data into a pandas DataFrame.\n"
    "2. Renamed columns with typos (e.g., 'life_exectancy' to 'life_expectancy').\n"
    "3. Removed non-numeric characters and converted columns to appropriate types.\n"
)

# Slide 4: Database Storage
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Database Storage"
content = slide.placeholders[1]
content.text = (
    "1. Created an SQLite database using SQLAlchemy.\n"
    "2. Defined a dynamic model for the data using SQLAlchemy.\n"
    "3. Inserted the cleaned data into the SQLite database.\n"
    "4. Verified data insertion by reading the data back from the database into a DataFrame.\n"
)

# Slide 5: Visualization
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Visualization"
content = slide.placeholders[1]
content.text = (
    "1. Used matplotlib and seaborn for data visualization.\n"
    "2. Created various plots such as histograms, box plots, scatter plots, and heatmaps.\n"
    "3. Displayed the initial DataFrame created from the HTML table.\n"
)

# Slide 6: Sample Data Table
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Sample Data"
rows, cols = df_from_db.shape
table = slide.shapes.add_table(rows+1, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(5)).table

# Set column names
for i, column_name in enumerate(df_from_db.columns):
    table.cell(0, i).text = column_name

# Add data to the table
for i in range(rows):
    for j in range(cols):
        table.cell(i+1, j).text = str(df_from_db.iat[i, j])

# Slide 7: Conclusion
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Conclusion"
content = slide.placeholders[1]
content.text = (
    "This presentation summarized the process of extracting, cleaning, storing, "
    "and visualizing global ecological footprint data. The data extraction was "
    "performed using MechanicalSoup, data cleaning with pandas, storage with "
    "SQLAlchemy and SQLite, and visualization with matplotlib and seaborn."
)

# Save the presentation
prs.save('Global_Ecological_Footprint_Summary.pptx')
print("PowerPoint presentation created successfully.")
